# MarkdownToPPTX.py

import os
import re
from typing import List, Tuple, Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Cm
from pptx.enum.text import PP_ALIGN


class MarkdownToPPTX:
    def __init__(self, template_path: Optional[str] = None):
        """
        Initialize the converter with a new presentation.
        Args:
            template_path (str, optional): PPTX template file path.
            if provided, creates a presentation based on the template;
            if not provided, creates a new blank presentation.
        """
        # default slide size
        default_width = Inches(13.333) # 16:9
        #default_width = Inches(10) # 4:3
        default_height = Inches(7.5)
        default_font_size = Pt(16)
        
        if template_path and os.path.exists(template_path):
            # load the presentation from the template file
            self.presentation = Presentation(template_path)
            # get the slide size from the template
            # Use the presentation's slide size directly
            self.presentation.slide_width = self.presentation.slide_width or default_width
            self.presentation.slide_height = self.presentation.slide_height or default_height
        elif template_path:
            # check whether the template file exists
            raise ValueError(f"Template file '{template_path}' does not exist.")
        elif os.path.exists('template.pptx'):
            self.presentation = Presentation('template.pptx')
            # get the slide size from the default template
            # Use the presentation's slide size directly
            self.presentation.slide_width = self.presentation.slide_width or default_width
            self.presentation.slide_height = self.presentation.slide_height or default_height
        else:
            # creates a new blank presentation
            self.presentation = Presentation()
            # set the default slide size
            self.presentation.slide_width = default_width
            self.presentation.slide_height = default_height
        
    def parse_markdown(self, markdown_text: str) -> List[dict]:
        """
        Parse markdown text and convert to structured slides data.
        
        Args:
            markdown_text (str): The markdown content to parse
            
        Returns:
            list: List of slide dictionaries containing title and content
        """
        # Split the markdown by slide separators
        slide_sections = re.split(r'\n---+\s*\n', markdown_text.strip())
        slides_data = []
        
        for section in slide_sections:
            lines = section.strip().split('\n')
            current_slide = None
            
            i = 0
            while i < len(lines):
                line = lines[i].rstrip()  # Keep leading spaces for indentation
                
                # Handle headers (slide titles and content headers)
                header_match = re.match(r'^(#{1,6})\s+(.*)', line)
                if header_match:
                    level = len(header_match.group(1))
                    title = header_match.group(2).strip()
                    
                    # # and ## create new slides
                    if level <= 2:
                        if current_slide:
                            slides_data.append(current_slide)
                        current_slide = {
                            'title': title,
                            'headers': [],
                            'content': []
                        }
                    # ###, ####, etc. are content headers
                    else:
                        if current_slide is not None:
                            current_slide['headers'].append({
                                'type': 'header',
                                'level': level,
                                'text': title
                            })
                        # If no slide exists yet, create one with default title
                        elif level > 2:
                            current_slide = {
                                'title': 'Content',
                                'headers': [{'type': 'header', 'level': level, 'text': title}],
                                'content': []
                            }
                # Handle bullet points
                elif re.match(r'^(\s*)(-|\*)\s+(.*)', line):
                    match = re.match(r'^(\s*)(-|\*)\s+(.*)', line)
                    if match:
                        spaces = match.group(1)
                        text = match.group(3)
                        
                        # Calculate indentation level (1 tab or 2/4 spaces = 1 level)
                        indent_level = 0
                        for char in spaces:
                            if char == '\t':
                                indent_level += 1
                            elif char == ' ':
                                # Group spaces into indentation levels
                                # (2 or 4 spaces = 1 level)
                                indent_level = (indent_level + 1) // 2
                        
                        if current_slide is not None:
                            current_slide['content'].append({
                                'type': 'bullet',
                                'level': indent_level,
                                'text': text
                            })
                        else:
                            # Create a slide if we encounter bullet before header
                            current_slide = {
                                'title': 'Content',
                                'headers': [],
                                'content': [{
                                    'type': 'bullet',
                                    'level': indent_level,
                                    'text': text
                                }]
                            }
                # Handle tables - check for table pattern
                elif '|' in line and line.strip() and current_slide is not None:
                    # Look ahead to see if this is actually a table
                    table_lines = []
                    j = i
                    
                    # Collect potential table lines
                    while j < len(lines) and '|' in lines[j]:
                        table_lines.append(lines[j])
                        j += 1
                        
                        # Check if next line is a separator line (:--- format)
                        if j < len(lines) and re.match(r'^\s*\|?(\s*:?-+:?\s*\|)+\s*:?-+:?\s*\|?\s*$', lines[j]):
                            table_lines.append(lines[j])  # Add separator line
                            j += 1
                            
                            # Collect additional data rows
                            while j < len(lines) and '|' in lines[j]:
                                table_lines.append(lines[j])
                                j += 1
                            break
                    
                    # If we found a valid table
                    if len(table_lines) >= 2:  # Need at least header + separator
                        current_slide['content'].append({
                            'type': 'table',
                            'lines': table_lines
                        })
                        i = j - 1  # Skip processed lines
                    else:
                        # Not a table, treat as regular text
                        self._handle_regular_text(line, current_slide)
                        
                elif '|' in line and line.strip() and current_slide is None:
                    # Create a slide if we encounter potential table before header
                    current_slide = {
                        'title': 'Content',
                        'headers': [],
                        'content': []
                    }
                    # Process the line again in the next iteration
                    continue
                # Handle regular paragraphs
                elif line.strip() and current_slide is not None:
                    self._handle_regular_text(line, current_slide)
                # Handle paragraphs when no slide has been created yet
                elif line.strip() and current_slide is None:
                    # Create a slide with a default title if we encounter content before a header
                    current_slide = {
                        'title': 'Content',
                        'headers': [],
                        'content': []
                    }
                    self._handle_regular_text(line, current_slide)
                    continue  # Continue to process the same line in case it's part of something else
                
                i += 1
            
            # Add the last slide of this section
            if current_slide:
                # Combine headers and content
                current_slide['content'] = current_slide['headers'] + current_slide['content']
                slides_data.append(current_slide)
                
        return slides_data

    def _handle_regular_text(self, line: str, current_slide: dict) -> None:
        """
        Helper method to handle regular text lines.
        
        Args:
            line (str): The line of text to handle
            current_slide (dict): The current slide being built
        """
        if line.strip():  # Skip empty lines
            # Check if this might be part of an unfinished table
            if not (current_slide.get('content') and 
                   current_slide['content'][-1]['type'] == 'table'):
                current_slide['content'].append({
                    'type': 'paragraph',
                    'text': line.strip()
                })

    def parse_table_data(self, table_lines: List[str]) -> List[List[str]]:
        """
        Parse table lines into structured data.
        
        Args:
            table_lines (List[str]): Lines containing table data
            
        Returns:
            List[List[str]]: Table data as list of rows, each row is a list of cell values
        """
        table_data = []
        in_header = True
        
        for line in table_lines:
            # Skip separator lines (:--- format)
            if re.match(r'^\s*\|?(\s*:?-+:?\s*\|)+\s*:?-+:?\s*\|?\s*$', line):
                in_header = False
                continue
                
            # Parse row data
            cells = [cell.strip() for cell in line.split('|')]
            
            # Remove empty cells at start and end
            if cells and cells[0] == '':
                cells = cells[1:]
            if cells and cells[-1] == '':
                cells = cells[:-1]
                
            if cells:
                table_data.append(cells)
                
        return table_data

    def remove_bold_formatting(self, text: str) -> Tuple[str, List[Tuple[int, int]]]:
        """
        Remove bold formatting (**text**) and return positions where bold should be applied.
        
        Args:
            text (str): Text that may contain bold formatting
            
        Returns:
            Tuple[str, List[Tuple[int, int]]]: Clean text and list of (start, end) positions for bold formatting
        """
        bold_positions = []
        clean_text = text
        offset = 0
        
        # Find all **bold** patterns
        for match in re.finditer(r'\*\*(.*?)\*\*', text):
            # Calculate positions in the cleaned text
            start = match.start() - offset
            end = start + len(match.group(1))
            bold_positions.append((start, end))
            
            # Remove the ** markers from the text
            clean_text = clean_text[:match.start()-offset] + match.group(1) + clean_text[match.end()-offset:]
            offset += 4  # 2 characters removed at start and 2 at end
            
        return clean_text, bold_positions

    def apply_text_formatting(self, paragraph, bold_ranges: List[Tuple[int, int]]) -> None:
        """
        Apply bold formatting to specific ranges in a paragraph.
        
        Args:
            paragraph: The paragraph object to format
            bold_ranges (List[Tuple[int, int]]): List of (start, end) positions for bold formatting
        """
        for start, end in bold_ranges:
            # Apply bold formatting to the specified range
            for run in paragraph.runs:
                # This is a simplified approach - in practice, you'd need to split runs appropriately
                if hasattr(run.font, 'bold'):
                    # We'll apply bold to entire runs for simplicity
                    run.font.bold = True

    def create_title_slide(self, title: str) -> object:
        """
        Create a title slide with the given title.
        
        Args:
            title (str): The title for the slide
            
        Returns:
            Slide: The created slide object
        """
        """ v0.1.0
        slide_layout = self.presentation.slide_layouts[0]  # Title Slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set title
        title_placeholder = slide.shapes.title
        clean_title, bold_positions = self.remove_bold_formatting(title)
        title_placeholder.text = clean_title
        
        # Apply bold formatting if needed
        if bold_positions and title_placeholder.text_frame.paragraphs:
            paragraph = title_placeholder.text_frame.paragraphs[0]
            self.apply_text_formatting(paragraph, bold_positions)
        """

        slide_layout = self.presentation.slide_layouts[0]  # Title Slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
    
        # 尝试获取标题占位符，如果没有则创建文本框
        title_placeholder = slide.shapes.title
    
        if title_placeholder is not None:
            # 使用现有的标题占位符
            clean_title, bold_positions = self.remove_bold_formatting(title)
            title_placeholder.text = clean_title
        
            # Apply bold formatting if needed
            if bold_positions and title_placeholder.text_frame.paragraphs:
                paragraph = title_placeholder.text_frame.paragraphs[0]
                self.apply_text_formatting(paragraph, bold_positions)
        else:
            # 如果没有标题占位符，则手动添加标题文本框
            clean_title, _ = self.remove_bold_formatting(title)
            title_box = slide.shapes.add_textbox(
                left=Cm(2), top=Cm(3), width=self.presentation.slide_width-Cm(4), height=Cm(3)
            )
            text_frame = title_box.text_frame
            text_frame.text = clean_title
            # 设置标题样式
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(44)
            p.font.bold = True


        return slide

    def create_content_slide(self, title: str, content: List[dict]) -> object:
        """
        Create a content slide with title and structured content.
        
        Args:
            title (str): The slide title
            content (list): List of content items with type, level, and text
            
        Returns:
            Slide: The created slide object
        """
        slide_layout = self.presentation.slide_layouts[1]  # Title and Content layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # 设置幻灯片标题
        title_placeholder = slide.shapes.title
        if title_placeholder is not None:
            clean_title, bold_positions = self.remove_bold_formatting(title)
            title_placeholder.text = clean_title
            
            # Apply bold formatting to title if needed
            if bold_positions and title_placeholder.text_frame.paragraphs:
                paragraph = title_placeholder.text_frame.paragraphs[0]
                self.apply_text_formatting(paragraph, bold_positions)
        else:
            # 如果没有标题占位符，则手动添加标题
            clean_title, _ = self.remove_bold_formatting(title)
            title_box = slide.shapes.add_textbox(
                left=Cm(1), top=Cm(0.5), width=self.presentation.slide_width-Cm(2), height=Cm(1.5)
            )
            text_frame = title_box.text_frame
            text_frame.text = clean_title
            p = text_frame.paragraphs[0]
            p.font.size = Pt(28)
            p.font.bold = True
        
        # 获取内容占位符
        content_placeholder = None
        try:
            content_placeholder = slide.placeholders[1]
        except KeyError:
            pass
        
        # 如果没有内容占位符，则创建一个新的文本框
        if content_placeholder is not None:
            content_text_frame = content_placeholder.text_frame
            content_text_frame.clear()  # 清除任何现有内容
        else:
            content_box = slide.shapes.add_textbox(
                left=Cm(1), top=Cm(2), width=self.presentation.slide_width-Cm(2), height=self.presentation.slide_height-Cm(3)
            )
            content_text_frame = content_box.text_frame
        
        # Position for tables
        current_top = Inches(1.5)
        left_margin = Inches(1)
        content_width = Inches(8)
        
        # 添加内容
        for item in content:
            if item['type'] == 'header':
                # 添加标题
                p = content_text_frame.paragraphs[0] if len(content_text_frame.paragraphs) == 1 and not content_text_frame.paragraphs[0].text else content_text_frame.add_paragraph()
                clean_text, bold_positions = self.remove_bold_formatting(item['text'])
                p.text = clean_text
                p.level = max(0, item['level'] - 3)  # 调整级别以适应演示文稿 (### = level 0)
                # 根据标题级别设置字体大小
                font_size = max(16, 28 - (item['level'] - 3) * 2)  # 最小尺寸 16pt
                p.font.size = Pt(font_size)
                p.font.bold = True
                
                # 如有必要，应用文本中的粗体格式
                if bold_positions:
                    self.apply_text_formatting(p, bold_positions)
            elif item['type'] == 'bullet':
                # 添加项目符号
                p = content_text_frame.paragraphs[0] if len(content_text_frame.paragraphs) == 1 and not content_text_frame.paragraphs[0].text else content_text_frame.add_paragraph()
                clean_text, bold_positions = self.remove_bold_formatting(item['text'])
                p.text = clean_text
                p.level = item['level']
                p.font.size = Pt(18)
                
                # 如有必要，应用文本中的粗体格式
                if bold_positions:
                    self.apply_text_formatting(p, bold_positions)
            elif item['type'] == 'paragraph':
                # 添加段落
                p = content_text_frame.paragraphs[0] if len(content_text_frame.paragraphs) == 1 and not content_text_frame.paragraphs[0].text else content_text_frame.add_paragraph()
                clean_text, bold_positions = self.remove_bold_formatting(item['text'])
                p.text = clean_text
                p.font.size = Pt(16)
                
                # 如有必要，应用文本中的粗体格式
                if bold_positions:
                    self.apply_text_formatting(p, bold_positions)
            elif item['type'] == 'table':
                # 解析并添加表格
                table_data = self.parse_table_data(item['lines'])
                if table_data:
                    rows = len(table_data)
                    cols = max(len(row) for row in table_data) if table_data else 0
                    
                    if rows > 0 and cols > 0:
                        # 创建表格
                        table_height = min(Inches(4), Inches(0.3 * rows))
                        table = slide.shapes.add_table(
                            rows, cols,
                            left_margin,
                            current_top,
                            content_width,
                            table_height
                        ).table
                        
                        # 填充表格
                        for row_idx, row_data in enumerate(table_data):
                            for col_idx, cell_data in enumerate(row_data):
                                if row_idx < rows and col_idx < cols:
                                    clean_cell_data, _ = self.remove_bold_formatting(cell_data)
                                    table.cell(row_idx, col_idx).text = clean_cell_data
                        
                        # 设置表格样式
                        for row_idx in range(rows):
                            for col_idx in range(cols):
                                cell = table.cell(row_idx, col_idx)
                                for paragraph in cell.text_frame.paragraphs:
                                    paragraph.font.size = Pt(12)
                                    if row_idx == 0:  # 标题行
                                        paragraph.font.bold = True
                        
                        current_top += table_height + Inches(0.2)
        
        return slide




    def get_unique_output_path(self, base_path: str) -> str:
        """
        Generate a unique output path by adding a counter if file already exists.
        
        Args:
            base_path (str): The base output path
            
        Returns:
            str: A unique output path
        """
        path = Path(base_path)
        counter = 1
        new_path = path
        
        while new_path.exists():
            # Split filename and extension
            name = path.stem
            ext = path.suffix
            # Create new filename with counter
            new_name = f"{name}({counter}){ext}"
            new_path = path.parent / new_name
            counter += 1
            
        return str(new_path)

    def convert(self, input_file_path: str, output_dir: str = "./output") -> None:
        """
        Convert markdown file to PPTX presentation.
        
        Args:
            input_file_path (str): Path to the input markdown file
            output_dir (str): Directory to save the output presentation
        """
        # Ensure output directory exists
        os.makedirs(output_dir, exist_ok=True)
        
        # Read markdown file
        try:
            with open(input_file_path, 'r', encoding='utf-8') as f:
                markdown_content = f.read()
        except FileNotFoundError:
            print(f"Error: Input file '{input_file_path}' not found.")
            return
        except Exception as e:
            print(f"Error reading file: {e}")
            return
        
        # Parse markdown content
        slides_data = self.parse_markdown(markdown_content)
        
        if not slides_data:
            print("Warning: No valid slide data found in markdown file.")
            return
        
        # Create slides
        first_header_slide = True
        for i, slide_data in enumerate(slides_data):
            # Check if this is the first # header to create a title slide
            if first_header_slide and slides_data[0]['title']:
                # Create title slide for the first main header
                self.create_title_slide(slide_data['title'])
                first_header_slide = False
                
                # If this slide has content, create a content slide too
                if slide_data['content']:
                    self.create_content_slide(slide_data['title'], slide_data['content'])
            else:
                self.create_content_slide(slide_data['title'], slide_data['content'])
        
        # Generate unique output path
        output_filename = "output.pptx"
        output_path = os.path.join(output_dir, output_filename)
        unique_output_path = self.get_unique_output_path(output_path)
        
        # Save presentation
        try:
            self.presentation.save(unique_output_path)
            print(f"Presentation saved to {unique_output_path}")
        except Exception as e:
            print(f"Error saving presentation: {e}")

def main():
    """
    Main function to run the markdown to PPTX converter.
    """
    # Create converter instance
    converter = MarkdownToPPTX("./assets/templates/template.pptx")
    
    # Define input and output paths
    input_file = r"./data/raw/sample.md"
    output_dir = r"./data/processed"
    
    # Convert markdown to PPTX
    converter.convert(input_file, output_dir)

if __name__ == "__main__":
    main()